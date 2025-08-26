"""PowerPoint file analyzer"""

from typing import Dict, Any, List
from .base import AnalyzerBase
from pptx import Presentation
import re


class PPTAnalyzer(AnalyzerBase):
    """Analyzer for PowerPoint files (.pptx)"""
    
    def analyze(self, file_path: str) -> Dict[str, Any]:
        """Analyze PowerPoint file and extract information"""
        self.validate_file(file_path)
        
        prs = Presentation(file_path)
        
        slides_data = []
        all_text = []
        
        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_info = self._analyze_slide(slide, slide_num)
            slides_data.append(slide_info)
            all_text.extend(slide_info['text_content'])
        
        # Extract tables and charts
        tables = self._extract_tables(prs)
        images_count = self._count_images(prs)
        
        return {
            "file_info": self.get_file_info(file_path),
            "presentation_info": {
                "slide_count": len(prs.slides),
                "title": self._get_presentation_title(prs),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height
            },
            "slides": slides_data,
            "statistics": {
                "total_slides": len(prs.slides),
                "total_text_boxes": sum(s['text_boxes'] for s in slides_data),
                "total_words": len(' '.join(all_text).split()),
                "tables_count": len(tables),
                "images_count": images_count
            },
            "tables": tables,
            "full_text": '\n\n'.join(all_text),
            "outline": self._generate_outline(slides_data)
        }
    
    def extract_text(self, file_path: str) -> str:
        """Extract all text from PowerPoint"""
        self.validate_file(file_path)
        
        prs = Presentation(file_path)
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"=== Slide {slide_num} ==="]
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            if len(slide_text) > 1:  # Has content beyond slide marker
                all_text.append('\n'.join(slide_text))
        
        return '\n\n'.join(all_text)
    
    def _analyze_slide(self, slide, slide_num: int) -> Dict[str, Any]:
        """Analyze individual slide"""
        text_content = []
        text_boxes = 0
        has_title = False
        title = ""
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_content.append(text)
                    text_boxes += 1
                    
                    # Check if it's a title
                    if shape == slide.shapes.title:
                        has_title = True
                        title = text
        
        # Try to detect title from layout if not found
        if not title and text_content:
            title = text_content[0][:50]  # First 50 chars as title
        
        return {
            "slide_number": slide_num,
            "title": title,
            "has_title": has_title,
            "text_boxes": text_boxes,
            "text_content": text_content,
            "layout": slide.slide_layout.name if slide.slide_layout else "Custom"
        }
    
    def _extract_tables(self, presentation: Presentation) -> List[Dict[str, Any]]:
        """Extract all tables from presentation"""
        tables = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            for shape in slide.shapes:
                if shape.has_table:
                    table_data = []
                    table = shape.table
                    
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(row_data)
                    
                    tables.append({
                        "slide": slide_num,
                        "rows": len(table.rows),
                        "columns": len(table.columns),
                        "data": table_data
                    })
        
        return tables
    
    def _count_images(self, presentation: Presentation) -> int:
        """Count images in presentation"""
        count = 0
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture type
                    count += 1
        return count
    
    def _get_presentation_title(self, presentation: Presentation) -> str:
        """Try to get presentation title from first slide"""
        if presentation.slides:
            first_slide = presentation.slides[0]
            if first_slide.shapes.title:
                return first_slide.shapes.title.text.strip()
            
            # Try to find largest text on first slide
            for shape in first_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    return shape.text.strip()[:100]
        
        return "Untitled Presentation"
    
    def _generate_outline(self, slides_data: List[Dict[str, Any]]) -> List[str]:
        """Generate presentation outline from slide titles"""
        outline = []
        for slide in slides_data:
            if slide['title']:
                outline.append(f"{slide['slide_number']}. {slide['title']}")
        return outline